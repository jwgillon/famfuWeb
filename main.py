import copy
import io
import logging
import os
import traceback
from datetime import datetime
from typing import Any

from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import HTMLResponse, Response
from fastapi.staticfiles import StaticFiles
from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
import httpx
from pydantic import BaseModel

logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s")
log = logging.getLogger(__name__)

app = FastAPI(title="Family Feud Generator API")

@app.get("/health")
def health():
    log.info("Health check hit")
    return {"status": "ok"}

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["POST", "GET"],
    allow_headers=["*"],
)

HERE = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_FULL   = os.path.join(HERE, "template.pptm")
TEMPLATE_SIMPLE = os.path.join(HERE, "template_simple.pptx")
GEMINI_MODEL = "gemini-2.5-flash"
MASTER_API_KEY = os.environ.get("GEMINI_API_KEY", "")


class GenerateRequest(BaseModel):
    game_data: dict[str, Any]
    theme: str = "Game"
    api_key: str = ""
    version: str = "full"


# ── Shape text helpers ────────────────────────────────────────────────────────

def set_shape_text(shape_collection, shape_name: str, text: str) -> bool:
    """Recursively search shapes (including inside groups) and set text."""
    for shape in shape_collection:
        if shape.shape_type == 6:  # GROUP — recurse
            if set_shape_text(shape.shapes, shape_name, text):
                return True
            continue
        if shape.name != shape_name:
            continue
        if not shape.has_text_frame:
            log.warning("Shape %s has no text frame", shape_name)
            return False
        tf = shape.text_frame
        if not tf.paragraphs:
            return False
        para = tf.paragraphs[0]
        if not para.runs:
            para.add_run().text = text
        else:
            para.runs[0].text = text
            for run in para.runs[1:]:
                run.text = ""
        return True
    return False


def set_multiline_shape(slide, shape_name: str, text: str) -> bool:
    """Replace all paragraphs in a shape with new multiline text.
    Used for print answer sheet columns which have many paragraphs."""
    for shape in slide.shapes:
        if shape.name != shape_name:
            continue
        if not shape.has_text_frame:
            return False
        tf = shape.text_frame
        txBody = tf._txBody

        # Save font properties from first run
        rPr_copy = None
        if tf.paragraphs and tf.paragraphs[0].runs:
            orig_rPr = tf.paragraphs[0].runs[0]._r.find(qn('a:rPr'))
            if orig_rPr is not None:
                rPr_copy = copy.deepcopy(orig_rPr)

        # Remove ALL existing paragraphs
        for p in txBody.findall(qn('a:p')):
            txBody.remove(p)

        # Insert one paragraph per line
        for line in text.split('\n'):
            p_elem = etree.fromstring(
                f'<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                f'<a:r><a:t>{line}</a:t></a:r></a:p>'
            )
            if rPr_copy is not None:
                r_elem = p_elem.find(qn('a:r'))
                r_elem.insert(0, copy.deepcopy(rPr_copy))
            txBody.append(p_elem)

        return True
    return False


def set_slide_shape(slide, shape_name: str, text: str) -> bool:
    return set_shape_text(slide.shapes, shape_name, text)


def set_all_slides(prs, shape_name: str, text: str) -> bool:
    for slide in prs.slides:
        if set_slide_shape(slide, shape_name, text):
            return True
    return False


# ── Shape visibility helper ───────────────────────────────────────────────────

def set_shape_visible(slide, shape_name: str, visible: bool) -> bool:
    """Set visibility on a top-level shape and all its children via XML."""
    for shape in slide.shapes:
        if shape.name != shape_name:
            continue
        # Set visibility on the shape itself
        cNvPr = shape._element.find('.//' + qn('p:cNvPr'))
        if cNvPr is not None:
            if visible:
                cNvPr.attrib.pop('hidden', None)
            else:
                cNvPr.set('hidden', '1')
        # Also set visibility on ALL children if it's a group
        if shape.shape_type == 6:  # msoGroup
            for child in shape.shapes:
                child_cNvPr = child._element.find('.//' + qn('p:cNvPr'))
                if child_cNvPr is not None:
                    if visible:
                        child_cNvPr.attrib.pop('hidden', None)
                    else:
                        child_cNvPr.set('hidden', '1')
        return True
    log.warning("Shape not found for visibility: %s", shape_name)
    return False


# ── Print answer sheet helper ─────────────────────────────────────────────────

def format_question_block(n: int, question: str, answers: list, points: list) -> str:
    lines = [f"QUESTION {n}: {question.upper()}"]
    for i, (ans, pts) in enumerate(zip(answers, points), 1):
        if str(ans).upper().strip() == "BLANK":
            continue
        prefix = f"{i}) {ans}"
        padded = prefix.ljust(20, '.')
        lines.append(f"{padded}({pts})")
    return "\n".join(lines)


# ── API endpoints ─────────────────────────────────────────────────────────────

@app.get("/has-master-key")
async def has_master_key():
    return {"available": bool(MASTER_API_KEY)}


class GeminiRequest(BaseModel):
    prompt: str


@app.post("/gemini")
async def call_gemini(req: GeminiRequest):
    if not MASTER_API_KEY:
        raise HTTPException(400, "No master API key configured on server")
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{GEMINI_MODEL}:generateContent?key={MASTER_API_KEY}"
    payload = {
        "contents": [{"parts": [{"text": req.prompt}]}],
        "generationConfig": {"temperature": 0.7, "maxOutputTokens": 8192}
    }
    async with httpx.AsyncClient(timeout=120) as client:
        res = await client.post(url, json=payload)
    if res.status_code != 200:
        raise HTTPException(502, f"Gemini API error: {res.text[:200]}")
    data = res.json()
    text = data.get("candidates", [{}])[0].get("content", {}).get("parts", [{}])[0].get("text", "")
    if not text:
        raise HTTPException(502, "Empty response from Gemini")
    log.info("Gemini via master key returned %d chars", len(text))
    return {"text": text}


@app.post("/generate")
async def generate(req: GenerateRequest):
    data = req.game_data

    if "questions" not in data:
        raise HTTPException(400, "Missing field: questions")

    questions = data["questions"]
    if len(questions) != 10:
        raise HTTPException(400, f"Expected 10 questions, got {len(questions)}")

    for i, q in enumerate(questions):
        if "q" not in q or "answers" not in q or "points" not in q:
            raise HTTPException(400, f"Question {i+1} missing q/answers/points")
        if len(q["answers"]) != 8:
            raise HTTPException(400, f"Question {i+1} needs 8 answers")
        if len(q["points"]) != 8:
            raise HTTPException(400, f"Question {i+1} needs 8 points")

    version = req.version.strip().lower()

    if version == "simple":
        template_path = TEMPLATE_SIMPLE
        media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    else:
        template_path = TEMPLATE_FULL
        media_type = "application/vnd.ms-powerpoint.presentation.macroEnabled.12"

    if not os.path.exists(template_path):
        raise HTTPException(500, f"Template not found: {template_path}")

    log.info("Generating Family Feud. Theme: %s, Version: %s", req.theme, version)

    try:
        prs = Presentation(template_path)
        log.info("Template loaded: %s — Slides: %d", template_path, len(prs.slides))

        def find_slide_by_name(name):
            for s in prs.slides:
                if s.name == name:
                    return s
            raise HTTPException(500, f"Data slide not found: '{name}'")

        data_slide  = find_slide_by_name("DataSlide_FamilyFeud")
        print_slide = find_slide_by_name("DataSlide_PrintCopy")

        blocks = []

        for q_idx, q in enumerate(questions):
            n = q_idx + 1
            question_text = str(q["q"])
            answers = q["answers"]
            points  = [str(p) for p in q["points"]]

            log.info("Q%d: '%s'", n, question_text[:50])

            # ── Data slide ──
            set_slide_shape(data_slide, f"FF_Q{n}", question_text)
            for a_idx in range(8):
                set_slide_shape(data_slide, f"FF_Q{n}_A{a_idx+1}", str(answers[a_idx]))
                set_slide_shape(data_slide, f"FF_Q{n}_P{a_idx+1}", str(points[a_idx]))

            # ── Question text on all 3 game slides ──
            for shape_name in [f"FF_QT{n}", f"FF_QQ{n}", f"FF_QTR{n}"]:
                set_all_slides(prs, shape_name, question_text)

            # ── Board slide (covers + answer groups) ──
            base = 3 if version == "simple" else 4
            board_slide_idx = base + (q_idx * 3) + 1
            if board_slide_idx < len(prs.slides):
                board_slide = prs.slides[board_slide_idx]
                for a_idx in range(8):
                    a_num    = a_idx + 1
                    is_blank = str(answers[a_idx]).upper().strip() == "BLANK"

                    # Write text
                    ok_a = set_slide_shape(board_slide, f"FF_Q{n}_AT{a_num}", str(answers[a_idx]))
                    ok_p = set_slide_shape(board_slide, f"FF_Q{n}_AP{a_num}", str(points[a_idx]))
                    if not ok_a or not ok_p:
                        log.warning("MISS on board: Q%d A%d AT=%s AP=%s", n, a_num, ok_a, ok_p)
                    else:
                        log.info("  AT%d='%s' AP%d=%s", a_num, answers[a_idx], a_num, points[a_idx])

                    # Reset to visible first, then hide if BLANK
                    set_shape_visible(board_slide, f"FF_Q{n}_GRP{a_num}",   not is_blank)
                    set_shape_visible(board_slide, f"FF_Q{n}_Cover{a_num}", not is_blank)

                    if is_blank:
                        log.info("  Hiding board slot Q%d A%d (BLANK)", n, a_num)

            # ── Reveal slide (no covers, groups only) ──
            reveal_slide_idx = base + (q_idx * 3) + 2
            if reveal_slide_idx < len(prs.slides):
                reveal_slide = prs.slides[reveal_slide_idx]
                for a_idx in range(8):
                    a_num    = a_idx + 1
                    is_blank = str(answers[a_idx]).upper().strip() == "BLANK"
                    # Write text
                    ok_ra = set_slide_shape(reveal_slide, f"FF_Q{n}_AT{a_num}_R", str(answers[a_idx]))
                    ok_rp = set_slide_shape(reveal_slide, f"FF_Q{n}_AP{a_num}_R", str(points[a_idx]))
                    if not ok_ra or not ok_rp:
                        log.warning("MISS on reveal: Q%d A%d AT=%s AP=%s", n, a_num, ok_ra, ok_rp)

                    # Reset to visible first, then hide if BLANK
                    set_shape_visible(reveal_slide, f"FF_Q{n}_GRP{a_num}_R", not is_blank)

                    if is_blank:
                        log.info("  Hiding reveal slot Q%d A%d (BLANK)", n, a_num)

            blocks.append(format_question_block(n, question_text, answers, points))

        # ── Print answer sheet ──
        col1 = "\n\n".join(blocks[0:2])
        col2 = "\n\n".join(blocks[2:5])
        col3 = "\n\n".join(blocks[5:8])
        col4 = "\n\n".join(blocks[8:10])

        set_multiline_shape(print_slide, "PrintColumn1", col1)
        set_multiline_shape(print_slide, "PrintColumn2", col2)
        set_multiline_shape(print_slide, "PrintColumn3", col3)
        set_multiline_shape(print_slide, "PrintColumn4", col4)

        now = datetime.now()
        print_info = f"FAMILY FEUD:\nHOST ANSWER KEY\nDate: {now.strftime('%d %B %Y')}\nTime: {now.strftime('%I:%M %p')}"
        set_multiline_shape(print_slide, "PrintInfo", print_info)
        log.info("Print answer sheet updated")

        # ── Build file ──
        slug = "".join(c for c in req.theme[:20] if c.isalnum() or c in " -_").strip() or "Game"
        ext = "pptx" if version == "simple" else "pptm"
        filename = f"Family Feud - {slug}.{ext}"

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        file_bytes = buf.read()
        log.info("File built: %d bytes, version: %s", len(file_bytes), version)

        return Response(
            content=file_bytes,
            media_type=media_type,
            headers={"Content-Disposition": f'attachment; filename="{filename}"'},
        )

    except HTTPException:
        raise
    except Exception as e:
        log.error("Generation failed: %s", traceback.format_exc())
        raise HTTPException(500, f"Failed to generate file: {str(e)}")


static_dir = os.path.join(HERE, "static")
if os.path.exists(static_dir):
    app.mount("/static", StaticFiles(directory=static_dir), name="static")


@app.get("/", response_class=HTMLResponse)
async def root():
    html_path = os.path.join(HERE, "index.html")
    if not os.path.exists(html_path):
        raise HTTPException(404, "index.html not found")
    with open(html_path) as f:
        return f.read()
